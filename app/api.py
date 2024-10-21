import base64
import time
from flask import json, jsonify, render_template, request, redirect, url_for
from pypinyin import pinyin
from pptx import Presentation
from PyPDF2 import PdfReader
from docx import Document
import pandas as pd

from PIL import Image
from app import app
from gtts import gTTS
from fuzzywuzzy import fuzz

import os
import cv2
import easyocr
import io
import gc

words = []
lang = None

def _submit_words(words_list, language):
    global words, lang
    words = [word.strip() for word in words_list if word.strip()]
    lang = language
    return 'Words submitted'

def _easyocr(image_file):
    image_path = './app/static/images/uploaded_image.jpg'
    image_file.save(image_path)
    img = cv2.imread(image_path)
    blur = cv2.GaussianBlur(img,(5,5),0)
    reader = easyocr.Reader(['ch_sim', 'en'])
    result = reader.readtext(blur)
    results = []
    for (bbox, text, _) in result:
        (_, _, _, _) = bbox
        results.append(text)
    del result
    gc.collect()
    os.remove(image_path)
    return jsonify({'text': '\n'.join(results)})

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/submit_words', methods=['POST'])
def submit_words_route():
    words_list = request.form['words'].split('\n')
    language = request.form['lang']
    _submit_words(words_list, language)
    return redirect(url_for('dictation_page'))

@app.route('/dictation_page.html')
def dictation_page():
    global words
    return render_template('dictation.html', words=words)

@app.route('/play_word/<string:word>/<string:lang>', methods=['GET'])
def play_word(word, lang):
    print("User typed: ", word, lang)
    try:
        tts = gTTS(text=word, lang=lang, slow=False)
        tts.save(f"app/static/audio/{word}_{lang}.mp3")
        return jsonify({'status': 'success'})
    except Exception as e:
        print(e)
        return jsonify({'status': 'error'}), 500

@app.route('/delete_audios.json', methods=['GET'])
def delete_audios():
    audio_dir = 'app/static/audio'
    if os.path.exists(audio_dir):
        for file in os.listdir(audio_dir):
            os.remove(os.path.join(audio_dir, file))
    return jsonify({'status': 'success'})

@app.route('/get_langs.json', methods=['GET'])
def get_langs():
    return jsonify({'langs': lang})

@app.route('/dict.html')
def dict():
    return render_template('dict.html',words=words)

@app.route('/convert_pinyin', methods=['GET'])
def convert_pinyin():
    data = request.args.get('data')
    pinyin_result = pinyin(data)
    pinyin_text = ' '.join([''.join(item) for item in pinyin_result])
    return jsonify({'pinyin': pinyin_text})

@app.route('/ocr_v1', methods=['POST'])
def easyocr_api():
    if 'image' not in request.files:
        return jsonify({'error': 'No image file provided\n没有提供图像文件'}), 400

    image_file = request.files['image']
    return _easyocr(image_file)

@app.route('/similarities_v1', methods=['POST'])
def get_similarities():
    data = request.get_json()
    results = {}

    for key, value in data.items():
        try:
            safe_key = key.encode('utf-8').hex() 
            image_string = value.split(",")[1]
            image_bytes = base64.b64decode(image_string)
            img = Image.open(io.BytesIO(image_bytes))
            img = img.convert('RGBA')
            background = Image.new('RGBA', img.size, (255, 255, 255))
            img = Image.alpha_composite(background, img)
            
            # Save the image to a temporary location
            image_path = f'./app/static/images/{safe_key}.png'
            img.save(image_path, format='PNG')

        except Exception as e:
            results[key] = {'error': f'Failed to decode base64 image: {str(e)}'}
            continue

        try:
            # Process the image with OCR
            reader = easyocr.Reader(['ch_sim', 'en'])
            ocr_results = reader.readtext(image_path)

            # Clean up: remove the image after processing (optional)
            os.remove(image_path)
            # Extract the recognized text from the OCR results
            if len(ocr_results) > 0:
                extracted_text = ' '.join([text[1] for text in ocr_results]).lower().replace(' ', '')
            else:
                results[key] = {'ocr_text': '', 'similarity': 0.0}
                continue

            # Calculate similarity between the key (word) and the OCR extracted text
            similarity = fuzz.ratio(extracted_text, key.lower())

            # Store the result
            results[key] = {'ocr_text': extracted_text, 'similarity': similarity}

        except Exception as e:
            results[key] = {'error': f'OCR processing error: {str(e)}'}

    return jsonify(results)

# @app.route('/ocr', methods=['POST'])
# def ocr_page():
#     if 'image' not in request.files:
#         return jsonify({'error': 'No image file provided\n没有提供图像文件'}), 400

#     image_file = request.files['image']
#     image_data = base64.b64encode(image_file.read()).decode('utf-8')
#     image_url = f"data:image/png;base64,{image_data}"
#     headers = {'Content-Type': 'application/json'}
#     data = {
#         "data": [image_url, ["en", "ch_sim"]],
#         "session_hash": "k65qy1668ak",
#         "action": "predict"
#     }
#     response = requests.post('https://tomofi-easyocr.hf.space/api/queue/push/', headers=headers, data=json.dumps(data))
#     new_data = {
#         "hash":response.json()['hash']
#     }
#     while True:
#         response = requests.post('https://tomofi-easyocr.hf.space/api/queue/status/', headers=headers, data=json.dumps(new_data))
#         if response.json()['status'] == "COMPLETE":
#             break
#         else:
#             time.sleep(2)
#     response = response.json()['data']['data'][1]['data']
#     extracted_texts = [text[0] for text in response]
#     return jsonify({'text': '\n'.join(extracted_texts)})

# @app.route('/similarities', methods=['POST'])
# def get_similarities_page():
#     data = request.get_json()
#     results = {}

#     for key, value in data.items():
#         try:
#             headers = {'Content-Type': 'application/json'}
#             data = {
#                 "data": [value, ["en", "ch_sim"]],
#                 "session_hash": "k65qy1668ak",
#                 "action": "predict"
#             }
#             response = requests.post('https://tomofi-easyocr.hf.space/api/queue/push/', headers=headers, data=json.dumps(data))
#             new_data = {
#                 "hash":response.json()['hash']
#             }
#             while True:
#                 response = requests.post('https://tomofi-easyocr.hf.space/api/queue/status/', headers=headers, data=json.dumps(new_data))
#                 if response.json()['status'] == "COMPLETE":
#                     break
#                 else:
#                     time.sleep(2)
#             response = response.json()['data']['data'][1]['data']
#             extracted_texts = [text[0] for text in response]
#             extracted_text = ' '.join(extracted_texts).lower().replace(' ', '')
#             similarity = fuzz.ratio(extracted_text, key.lower())

#             results[key] = {'ocr_text': extracted_text, 'similarity': similarity}
#             # results[key] = {'ocr_text': key, 'similarity': random.randint(60, 100)}

#         except Exception as e:
#             results[key] = {'error': f'OCR processing error: {str(e)}'}

#     return jsonify(results)

@app.route('/extract_docsfile_text', methods=['POST'])
def extract_docsfile_text():
    file = request.files['docsfile']
    if file.filename.endswith('.jpg') or file.filename.endswith('.png') or file.filename.endswith('.jpeg') or file.filename.endswith('.webp'):
        text_runs = []
        response = _easyocr(file)
        if response and response.is_json:
            text_data = response.get_json()
            if 'text' in text_data:
                text_runs.append(text_data['text'])
            else:
                return jsonify({'error': 'Unexpected response format from easyocr_api'}), 400
        else:
            return jsonify({'error': 'Error occurred while calling easyocr_api'}), 400
    elif file.filename.endswith('.pptx') or file.filename.endswith('.ppt'):
        prs = Presentation(file)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
                text_runs.append('\n')
        text_runs.pop(-1)
    elif file.filename.endswith('.pdf'):
        pdf = PdfReader(file)
        text_runs = []
        for page in pdf.pages:
            text_runs.append(page.extract_text())
        text_runs[-1] = text_runs[-1].replace('\n', '')
    elif file.filename.endswith(('.docx', '.doc', '.odt')):
        doc = Document(file)
        text_runs = []
        for para in doc.paragraphs:
            text_runs.append(para.text)
    elif file.filename.endswith('.txt'):
        text_runs = file.read().decode('utf-8').split('\n')
    elif file.filename.endswith(('.xlsx', '.xls', '.ods', '.csv')):
        text_runs = []
        if file.filename.endswith('.xlsx'):
            df = pd.read_excel(file)
            print("Loaded .xlsx file successfully.")
        elif file.filename.endswith('.xls'):
            df = pd.read_excel(file, engine='xlrd')
            print("Loaded .xls file successfully.")
        elif file.filename.endswith('.ods'):
            df = pd.read_excel(file, engine='odf')
            print("Loaded .ods file successfully.")
        elif file.filename.endswith('.csv'):
            df = pd.read_csv(file)
            print("Loaded .csv file successfully.")
        
        if isinstance(df,pd.DataFrame): 
            for sheet_name, sheet in df.items():
                print(f"Processing sheet: {type(sheet)} - {sheet_name}")
                if isinstance(sheet, pd.core.series.Series):
                    text_runs.extend([str(x) for x in sheet]) 
                else:
                    print(f"Unexpected format for sheet: {sheet_name}")
        elif isinstance(df, pd.DataFrame):  
            for index, row in df.iterrows():
                text_runs.append(str(row.to_dict()))  
        else:
                print("Unexpected format: df is neither a DataFrame nor a dict.")
                return {"error": "Unexpected format."}
    else:
        return jsonify({'error': 'File type not supported'})
    # return jsonify({'text': '\n'.join(text_runs)})
    return jsonify({'text': ''.join(text_runs)})
